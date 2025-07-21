import pandas as pd
from fpdf import FPDF
import pptx
from pptx.util import Inches
from google.cloud import storage
import os
import uuid

BUCKET_NAME = os.getenv("BUCKET_NAME", "your-gcs-bucket")
#BUCKET_NAME = "cxo-prism"

storage_client = storage.Client()


def upload_to_gcs(file_path):
    bucket = storage_client.bucket(BUCKET_NAME)
    destination_blob_name = f"reports/{uuid.uuid4().hex}/{os.path.basename(file_path)}"
    blob = bucket.blob(destination_blob_name)
    blob.upload_from_filename(file_path)
    return blob.public_url


def generate_pdf_report(df):
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", size=10)
    for index, row in df.iterrows():
        pdf.cell(200, 8, txt=str(row.to_dict()), ln=True)
    output_path = "report.pdf"
    pdf.output(output_path)
    return upload_to_gcs(output_path)


def generate_excel_report(df):
    output_path = "report.xlsx"
    df.to_excel(output_path, index=False)
    return upload_to_gcs(output_path)


def generate_ppt_report(df):
    prs = pptx.Presentation()
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    text_frame = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5)).text_frame
    for index, row in df.iterrows():
        text_frame.add_paragraph(str(row.to_dict()))
    output_path = "report.pptx"
    prs.save(output_path)
    return upload_to_gcs(output_path)